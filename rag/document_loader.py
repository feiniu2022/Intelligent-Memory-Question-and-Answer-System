import hashlib
import os
from datetime import datetime
from typing import Dict, List, Optional

from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import settings
from memory.vector_store import VectorStoreManager
from utils.logger import setup_logger

logger = setup_logger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt": "text", ".md": "markdown", ".markdown": "markdown",
    ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".bmp": "image", ".tiff": "image", ".webp": "image",
}

_ocr_engine = None


def _get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is not None:
        return _ocr_engine if _ocr_engine is not False else None
    try:
        from rapidocr_onnxruntime import RapidOCR
        _ocr_engine = RapidOCR()
        logger.info("RapidOCR 引擎初始化成功")
        return _ocr_engine
    except ImportError:
        pass
    try:
        from paddleocr import PaddleOCR
        _ocr_engine = PaddleOCR(lang="ch")
        logger.info("PaddleOCR 引擎初始化成功")
        return _ocr_engine
    except Exception as e:
        logger.warning("OCR 引擎初始化失败: %s", e)
        _ocr_engine = False
        return None


def _ocr_image(image_path: str) -> str:
    engine = _get_ocr_engine()
    if engine is None:
        return ""
    try:
        if 'RapidOCR' in type(engine).__name__:
            result, _ = engine(image_path)
            if not result:
                return ""
            lines = []
            for item in result:
                if isinstance(item, (list, tuple)) and len(item) >= 2:
                    lines.append(str(item[1]))
            text = "\n".join(lines)
            logger.info("RapidOCR 识别: %s (%d 字)", os.path.basename(image_path), len(text))
            return text
        else:
            result = engine.ocr(image_path, cls=True)
            if not result or not result[0]:
                return ""
            lines = []
            for line in result[0]:
                if line and len(line) >= 2:
                    lines.append(line[1][0])
            text = "\n".join(lines)
            logger.info("PaddleOCR 识别: %s (%d 字)", os.path.basename(image_path), len(text))
            return text
    except Exception as e:
        logger.error("OCR 识别失败 %s: %s", image_path, e)
        return ""


_splitter = RecursiveCharacterTextSplitter(
    separators=["\n\n", "\n", "\u3002", "\uff01", "\uff1f", "\uff1b", "\uff0c", " ", ""],
    chunk_size=settings.rag_chunk_size,
    chunk_overlap=settings.rag_chunk_overlap,
    length_function=len,
)


class DocumentLoader:
    """多格式文档加载器：TXT/MD/PDF/DOCX/PPTX/图片 → 向量库"""

    def __init__(self, persist_directory: Optional[str] = None, upload_dir: Optional[str] = None):
        persist_directory = persist_directory or str(settings.resolved_knowledge_db_dir)
        upload_dir = upload_dir or str(settings.resolved_upload_dir)
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)
        self.vector_manager = VectorStoreManager(
            persist_directory=persist_directory,
            collection_name="knowledge_base",
        )

    def _compute_file_hash(self, filepath: str) -> str:
        h = hashlib.md5()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()

    def _extract_text(self, filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        if ext in (".txt", ".md", ".markdown"):
            for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
                try:
                    with open(filepath, "r", encoding=enc) as f:
                        return f.read()
                except (UnicodeDecodeError, UnicodeError):
                    continue
            logger.warning("无法解码文件: %s", filepath)
            return ""
        if ext == ".pdf":
            return self._extract_pdf(filepath)
        if ext == ".docx":
            return self._extract_docx(filepath)
        if ext == ".pptx":
            return self._extract_pptx(filepath)
        if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
            return self._extract_image(filepath)
        logger.warning("不支持的格式: %s", ext)
        return ""

    def _extract_pdf(self, filepath: str) -> str:
        try:
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            text_parts = []
            has_text_page = False
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    has_text_page = True
                    text_parts.append(page_text.strip())

            if has_text_page:
                return "\n\n".join(text_parts)

            logger.info("PDF 无文字层，尝试 OCR: %s", filepath)
            return self._ocr_pdf_pages(filepath)
        except Exception as e:
            logger.error("PDF解析失败 %s: %s", filepath, e)
            return self._ocr_pdf_pages(filepath)

    def _extract_docx(self, filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.error("DOCX解析失败 %s: %s", filepath, e)
            return ""

    def _extract_pptx(self, filepath: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            slides = []
            for slide in prs.slides:
                texts = [para.text.strip() for shape in slide.shapes if shape.has_text_frame for para in shape.text_frame.paragraphs if para.text.strip()]
                if texts:
                    slides.append("\n".join(texts))
            return "\n\n---\n\n".join(slides)
        except Exception as e:
            logger.error("PPTX解析失败 %s: %s", filepath, e)
            return ""

    def _extract_image(self, filepath: str) -> str:
        text = _ocr_image(filepath)
        if not text:
            logger.warning("图片 OCR 无结果: %s", filepath)
        return text

    def _ocr_pdf_pages(self, filepath: str) -> str:
        engine = _get_ocr_engine()
        if engine is None:
            logger.warning("OCR 引擎不可用，跳过扫描PDF: %s", filepath)
            return ""
        try:
            from pdf2image import convert_from_path

            pages = convert_from_path(filepath, dpi=200)
            all_text = []
            tmp_dir = os.path.join(os.path.dirname(filepath), ".ocr_tmp")
            os.makedirs(tmp_dir, exist_ok=True)

            for i, page_img in enumerate(pages):
                img_path = os.path.join(tmp_dir, f"{os.path.basename(filepath)}_p{i}.png")
                page_img.save(img_path, "PNG")
                text = _ocr_image(img_path)
                if text:
                    all_text.append(f"--- 第{i+1}页 ---\n{text}")
                try:
                    os.remove(img_path)
                except OSError:
                    pass

            try:
                os.rmdir(tmp_dir)
            except OSError:
                pass

            if all_text:
                logger.info("OCR PDF 完成: %s (%d/%d 页有文字)", os.path.basename(filepath), len(all_text), len(pages))
            return "\n\n".join(all_text)
        except ImportError:
            logger.warning("pdf2image 未安装，扫描PDF OCR不可用。安装: pip install pdf2image")
            return ""
        except Exception as e:
            logger.error("OCR PDF 失败 %s: %s", filepath, e)
            return ""

    def load_file(self, filepath: str, user_id: str = "default_user") -> Dict:
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            return {"success": False, "error": f"不支持: {ext}", "chunks": 0}

        file_hash = self._compute_file_hash(filepath)
        existing = self.vector_manager.get_all(filter={"file_hash": file_hash})
        if existing and existing.get("ids"):
            logger.info("文件已存在: %s (跳过)", filename)
            return {"success": True, "error": None, "chunks": 0, "skipped": True, "filename": filename}

        old = self.vector_manager.get_all(filter={"$and": [{"filename": filename}, {"user_id": user_id}]})
        if old and old.get("ids"):
            self.vector_manager.delete_by_filter({"$and": [{"filename": filename}, {"user_id": user_id}]})
            logger.info("删除旧版本: %s", filename)

        logger.info("正在处理文件: %s ...", filename)
        content = self._extract_text(filepath)
        if not content.strip():
            return {"success": False, "error": "文件内容为空（图片/扫描件需安装 rapidocr-onnxruntime 识别）", "chunks": 0}

        chunks = _splitter.split_text(content)
        if not chunks:
            return {"success": False, "error": "分块为空", "chunks": 0}

        texts, metadatas, ids = [], [], []
        for i, chunk in enumerate(chunks):
            texts.append(chunk)
            metadatas.append({
                "filename": filename, "file_hash": file_hash, "chunk_index": i,
                "total_chunks": len(chunks), "source": SUPPORTED_EXTENSIONS.get(ext, "unknown"),
                "user_id": user_id, "timestamp": datetime.now().isoformat(),
            })
            ids.append(f"{file_hash}_{user_id}_{i}")

        self.vector_manager.add_texts(texts=texts, metadatas=metadatas, ids=ids)
        logger.info("已加载: %s (%d chunks, user=%s)", filename, len(chunks), user_id)
        return {"success": True, "error": None, "chunks": len(chunks), "filename": filename}

    def load_all_files(self, user_id: str = "default_user") -> int:
        txt_dir = str(settings.resolved_txt_data_dir)
        total = 0
        files = [f for f in os.listdir(txt_dir) if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS]
        if not files:
            logger.warning("目录为空: %s", txt_dir)
            return 0
        logger.info("加载文件 (%d 个)...", len(files))
        for fn in files:
            total += self.load_file(os.path.join(txt_dir, fn), user_id=user_id).get("chunks", 0)
        logger.info("加载完成: %d chunks", total)
        return total

    def search(self, query: str, k: int = 5, user_id: Optional[str] = None) -> List[Dict]:
        f = {"user_id": user_id} if user_id else None
        raw = self.vector_manager.similarity_search(query, k=k, filter=f)
        return [{"content": d.page_content, "metadata": d.metadata, "score": round(s, 4)} for d, s in raw]

    def list_files(self, user_id: Optional[str] = None) -> List[Dict]:
        f = {"user_id": user_id} if user_id else None
        results = self.vector_manager.get_all(filter=f)
        if not results or not results.get("metadatas"):
            return []
        seen, files = set(), []
        for m in results["metadatas"]:
            if m and m.get("filename") and m["filename"] not in seen:
                seen.add(m["filename"])
                files.append({"filename": m["filename"], "source": m.get("source", "unknown"), "timestamp": m.get("timestamp", ""), "total_chunks": m.get("total_chunks", 0)})
        return files

    def delete_file(self, filename: str, user_id: str = "default_user") -> bool:
        results = self.vector_manager.get_all(filter={"$and": [{"filename": filename}, {"user_id": user_id}]})
        if results and results.get("ids"):
            self.vector_manager.delete_by_filter({"$and": [{"filename": filename}, {"user_id": user_id}]})
            logger.info("已删除: %s (user=%s)", filename, user_id)
            return True
        return False

    def delete_all(self, user_id: str = "default_user") -> int:
        results = self.vector_manager.get_all(filter={"user_id": user_id})
        count = len(results.get("ids", [])) if results else 0
        if count > 0:
            self.vector_manager.delete_by_filter({"user_id": user_id})
        logger.info("已删除 %s 的所有文档 (%d chunks)", user_id, count)
        return count